"""Infinity Converter 6.0 high-value tool pack.

These tools intentionally use the project's local engines (pypdf/PyMuPDF/Pillow/
python-docx/openpyxl/Python stdlib) so no new cloud dependency is required.
"""
from __future__ import annotations

import base64, bz2, difflib, hashlib, html, io, json, lzma, mimetypes, re, tarfile, uuid, urllib.parse, zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import fitz
from PIL import Image, ImageChops, ImageEnhance, ImageFilter, ImageOps, ImageDraw
from docx import Document
from openpyxl import load_workbook


def _write_text(output: Path, text: str):
    output.write_text(text, encoding="utf-8")


def _ocr_lang(value: str) -> str:
    v = (value or "ar+eng").strip().lower().replace(" ", "")
    aliases = {"ar+en": "ara+eng", "en+ar": "eng+ara", "ar": "ara", "en": "eng", "ara": "ara", "eng": "eng", "ara+eng": "ara+eng", "eng+ara": "eng+ara"}
    return aliases.get(v, v or "ara+eng")


def _pdf_doc(path: Path):
    try:
        return fitz.open(path)
    except Exception as exc:
        raise ValueError("ملف PDF غير صالح أو تالف.") from exc


def pdf_to_docx(source: Path, output: Path):
    doc = Document()
    pdf = _pdf_doc(source)
    try:
        for i, page in enumerate(pdf):
            if i:
                doc.add_page_break()
            text = page.get_text("text").strip()
            doc.add_paragraph(text or "")
    finally:
        pdf.close()
    doc.save(output)


def pdf_to_markdown(source: Path, output: Path):
    pdf = _pdf_doc(source)
    try:
        chunks = []
        for i, page in enumerate(pdf, 1):
            text = page.get_text("text").strip()
            chunks.append(f"## Page {i}\n\n{text or '_No extractable text_'}")
        _write_text(output, "\n\n".join(chunks) + "\n")
    finally:
        pdf.close()


def pdf_compare(a: Path, b: Path, output: Path):
    pa, pb = _pdf_doc(a), _pdf_doc(b)
    try:
        ta = "\n".join(p.get_text("text") for p in pa).splitlines()
        tb = "\n".join(p.get_text("text") for p in pb).splitlines()
        diff = list(difflib.unified_diff(ta, tb, fromfile=a.name, tofile=b.name, lineterm=""))
        text = "\n".join(diff).strip()
        _write_text(output, (text + "\n") if text else "No differences found. The PDF texts are identical.\n")
    finally:
        pa.close(); pb.close()


def pdf_repair(source: Path, output: Path):
    pdf = _pdf_doc(source)
    try:
        pdf.save(output, garbage=4, deflate=True, clean=True)
    finally:
        pdf.close()


def pdf_image_extract(source: Path, output: Path):
    pdf = _pdf_doc(source)
    try:
        with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zf:
            seen = 0
            for page_num, page in enumerate(pdf, 1):
                for img_num, info in enumerate(page.get_images(full=True), 1):
                    xref = info[0]
                    try:
                        base = fitz.Pixmap(pdf, xref)
                        if base.alpha:
                            base = fitz.Pixmap(fitz.csRGB, base)
                        data = base.tobytes("png")
                        name = f"page-{page_num}-image-{img_num}.png"
                        zf.writestr(name, data)
                        seen += 1
                    except Exception:
                        continue
            if not seen:
                zf.writestr("README.txt", "No embedded raster images were found.\n")
    finally:
        pdf.close()


def pdf_links_report(source: Path, output: Path):
    pdf = _pdf_doc(source)
    rows=[]
    try:
        for i,p in enumerate(pdf,1):
            for link in p.get_links():
                rows.append({"page":i,"kind":link.get("kind"),"uri":link.get("uri"),"file":link.get("file")})
    finally: pdf.close()
    _write_text(output, json.dumps({"links":rows},ensure_ascii=False,indent=2))


def pdf_annotations_report(source: Path, output: Path):
    pdf = _pdf_doc(source); rows=[]
    try:
        for i,p in enumerate(pdf,1):
            annot=p.first_annot
            while annot:
                rows.append({"page":i,"type":annot.type[0],"name":annot.info.get("title"),"content":annot.info.get("content")})
                annot=annot.next
    finally: pdf.close()
    _write_text(output, json.dumps({"annotations":rows},ensure_ascii=False,indent=2))


def pdf_page_size_report(source: Path, output: Path):
    pdf = _pdf_doc(source); rows=[]
    try:
        for i,p in enumerate(pdf,1):
            r=p.rect; rows.append({"page":i,"width_points":round(r.width,2),"height_points":round(r.height,2),"width_mm":round(r.width*25.4/72,2),"height_mm":round(r.height*25.4/72,2)})
    finally: pdf.close()
    _write_text(output,json.dumps({"pages":rows},ensure_ascii=False,indent=2))


def pdf_redact(source: Path, output: Path, terms: str):
    needles=[x.strip() for x in re.split(r"[,\n]",terms or "") if x.strip()]
    if not needles: raise ValueError("اكتب كلمة أو أكثر للحذف الآمن، مفصولة بفواصل.")
    pdf=_pdf_doc(source)
    total=0
    try:
        for page in pdf:
            for needle in needles:
                for rect in page.search_for(needle):
                    page.add_redact_annot(rect, fill=(0,0,0)); total+=1
            page.apply_redactions()
        if total == 0:
            raise ValueError("لم يتم العثور على النص المطلوب تنقيحه.")
        pdf.save(output, garbage=4, deflate=True)
    finally: pdf.close()


def pdf_unlock(source: Path, output: Path, password: str):
    pdf=fitz.open(source)
    try:
        if pdf.needs_pass:
            if not pdf.authenticate(password or ""):
                raise ValueError("كلمة المرور غير صحيحة.")
        pdf.save(output, garbage=4, deflate=True)
    finally: pdf.close()


def image_upscale(source: Path, output: Path, factor: str = "2"):
    f=max(1,min(4,int(factor or 2)))
    with Image.open(source) as img:
        img=ImageOps.exif_transpose(img)
        out=img.resize((img.width*f,img.height*f), Image.Resampling.LANCZOS)
        out.save(output)


def image_blur(source: Path, output: Path, radius: str = "3"):
    r=max(0.1,min(50,float(radius or 3)))
    with Image.open(source) as img:
        ImageOps.exif_transpose(img).filter(ImageFilter.GaussianBlur(r)).save(output)


def image_pixelate(source: Path, output: Path, blocks: str = "32"):
    b=max(4,min(256,int(blocks or 32)))
    with Image.open(source) as img:
        img=ImageOps.exif_transpose(img).convert("RGB")
        w,h=img.size; small=img.resize((max(1,w//b),max(1,h//b)),Image.Resampling.BILINEAR)
        img=small.resize((w,h),Image.Resampling.NEAREST); img.save(output,quality=92)


def image_invert(source: Path, output: Path):
    with Image.open(source) as img:
        rgba=img.convert("RGBA"); rgb=ImageOps.invert(rgba.convert("RGB")); rgb.putalpha(rgba.getchannel("A")); rgb.save(output)


def image_posterize(source: Path, output: Path, bits: str = "4"):
    k=max(1,min(8,int(bits or 4)))
    with Image.open(source) as img: ImageOps.posterize(ImageOps.exif_transpose(img).convert("RGB"),k).save(output)


def image_palette(source: Path, output: Path, count: str = "8"):
    n=max(2,min(16,int(count or 8)))
    with Image.open(source) as img:
        pal=ImageOps.exif_transpose(img).convert("RGB").quantize(colors=n).convert("RGB")
        colors=pal.getcolors(maxcolors=1_000_000) or []
        top=sorted(colors,reverse=True)[:n]
        rows=[]
        for freq,color in top: rows.append({"rgb":list(color),"hex":"#%02x%02x%02x"%color,"pixels":freq})
        _write_text(output,json.dumps({"colors":rows},ensure_ascii=False,indent=2))


def image_watermark(source: Path, output: Path, text: str):
    text=(text or "INFINITY").strip()[:80]
    with Image.open(source) as img:
        img=ImageOps.exif_transpose(img).convert("RGBA")
        layer=Image.new("RGBA",img.size,(0,0,0,0)); d=ImageDraw.Draw(layer)
        size=max(18,min(72,img.width//18 if img.width else 24));
        try: from PIL import ImageFont; font=ImageFont.load_default(size=size)
        except Exception: font=None
        bbox=d.textbbox((0,0),text,font=font); tw,th=bbox[2]-bbox[0],bbox[3]-bbox[1]
        x=max(10,img.width-tw-24); y=max(10,img.height-th-24)
        d.rounded_rectangle((x-12,y-10,x+tw+12,y+th+10),radius=10,fill=(0,0,0,95)); d.text((x,y),text,fill=(255,255,255,190),font=font)
        Image.alpha_composite(img,layer).save(output)


def image_background_cleaner(source: Path, output: Path, tolerance: str = "24"):
    tol=max(1,min(120,int(tolerance or 24)))
    with Image.open(source) as img:
        img=ImageOps.exif_transpose(img).convert("RGBA")
        px=img.load();
        for y in range(img.height):
            for x in range(img.width):
                r,g,b,a=px[x,y]
                if min(r,g,b) > 245-tol:
                    px[x,y]=(r,g,b,0)
        img.save(output)


def image_auto_orient(source: Path, output: Path):
    with Image.open(source) as img: ImageOps.exif_transpose(img).save(output)


def image_round_corners(source: Path, output: Path, radius: str = "40"):
    r=max(1,min(500,int(radius or 40)))
    with Image.open(source) as img:
        img=ImageOps.exif_transpose(img).convert("RGBA")
        mask=Image.new("L",img.size,0); ImageDraw.Draw(mask).rounded_rectangle((0,0,img.width,img.height),radius=min(r,min(img.size)//2),fill=255)
        img.putalpha(mask); img.save(output)


def image_border(source: Path, output: Path, width: str = "24"):
    w=max(1,min(200,int(width or 24)))
    with Image.open(source) as img: ImageOps.expand(ImageOps.exif_transpose(img).convert("RGB"),border=w,fill="white").save(output,quality=95)


def docx_to_markdown(source: Path, output: Path):
    doc=Document(source); out=[]
    for p in doc.paragraphs:
        t=p.text.strip()
        if not t: out.append(""); continue
        style=p.style.name.lower() if p.style else ""
        out.append("# "+t if "title" in style else "## "+t if "heading 1" in style else "### "+t if "heading 2" in style else t)
    _write_text(output,"\n\n".join(out)+"\n")


def docx_table_csv(source: Path, output: Path):
    import csv
    doc=Document(source)
    rows=[]
    for table in doc.tables:
        for row in table.rows: rows.append([c.text.replace("\n"," ").strip() for c in row.cells])
    if not rows: raise ValueError("لم نجد جدولًا في ملف Word.")
    with output.open("w",encoding="utf-8",newline="") as f: csv.writer(f).writerows(rows)


def xlsx_to_html(source: Path, output: Path):
    wb=load_workbook(source,read_only=True,data_only=True); chunks=["<!doctype html><html><body>"]
    for ws in wb.worksheets:
        chunks.append(f"<h2>{html.escape(ws.title)}</h2><table><tbody>")
        for row in ws.iter_rows(values_only=True): chunks.append("<tr>"+"".join(f"<td>{html.escape(str(v)) if v is not None else ''}</td>" for v in row)+"</tr>")
        chunks.append("</tbody></table>")
    wb.close(); chunks.append("</body></html>"); _write_text(output,"\n".join(chunks))


def xlsx_summary(source: Path, output: Path):
    wb=load_workbook(source,read_only=True,data_only=True); sheets=[]
    try:
        for ws in wb.worksheets:
            rows=0; cols=0
            for row in ws.iter_rows(values_only=True):
                if any(v is not None for v in row): rows+=1; cols=max(cols,len(row))
            sheets.append({"sheet":ws.title,"rows":rows,"columns":cols})
    finally: wb.close()
    _write_text(output,json.dumps({"sheets":sheets},ensure_ascii=False,indent=2))


def csv_to_markdown(source: Path, output: Path):
    import csv
    with source.open("r",encoding="utf-8-sig",newline="") as f: rows=list(csv.reader(f))
    if not rows: raise ValueError("CSV فارغ.")
    width=len(rows[0]); rows=[r+['']*(width-len(r)) for r in rows]
    lines=["| "+" | ".join(rows[0])+" |","| "+" | ".join(["---"]*width)+" |"]+["| "+" | ".join(r)+" |" for r in rows[1:]]
    _write_text(output,"\n".join(lines)+"\n")


def csv_statistics(source: Path, output: Path):
    import pandas as pd
    df=pd.read_csv(source)
    report={"rows":int(len(df)),"columns":int(len(df.columns)),"columns_detail":[]}
    for c in df.columns:
        s=df[c]; detail={"name":str(c),"non_null":int(s.notna().sum()),"unique":int(s.nunique(dropna=True))}
        if pd.api.types.is_numeric_dtype(s): detail.update({"min":None if s.dropna().empty else float(s.min()),"max":None if s.dropna().empty else float(s.max()),"mean":None if s.dropna().empty else float(s.mean())})
        report["columns_detail"].append(detail)
    _write_text(output,json.dumps(report,ensure_ascii=False,indent=2,default=str))


def json_to_html(source: Path, output: Path):
    data=json.loads(source.read_text(encoding="utf-8"))
    _write_text(output,"<!doctype html><html><head><meta charset=\"utf-8\"><title>JSON</title></head><body><pre>"+html.escape(json.dumps(data,ensure_ascii=False,indent=2)) + "</pre></body></html>\n")


def html_to_text(source: Path, output: Path):
    raw=source.read_text(encoding="utf-8",errors="replace")
    text=re.sub(r"<script[\s\S]*?</script>|<style[\s\S]*?</style>"," ",raw,flags=re.I)
    text=re.sub(r"<[^>]+>"," ",text); text=html.unescape(re.sub(r"\s+"," ",text)).strip()
    _write_text(output,text+"\n")


def markdown_to_text(source: Path, output: Path):
    raw=source.read_text(encoding="utf-8")
    raw=re.sub(r"^\s{0,3}#{1,6}\s*", "", raw, flags=re.M); raw=re.sub(r"[*_`>~-]", "", raw)
    _write_text(output,raw)


def pptx_to_markdown(source: Path, output: Path):
    from pptx import Presentation
    prs=Presentation(source); out=[]
    for i,slide in enumerate(prs.slides,1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape,"text") and shape.text.strip(): out.append(shape.text.strip())
        out.append("")
    _write_text(output,"\n\n".join(out))


def _ocr_text_file(source: Path, lang: str = "ar+eng") -> str:
    import pytesseract
    from PIL import Image
    with Image.open(source) as img:
        return pytesseract.image_to_string(ImageOps.exif_transpose(img),lang=_ocr_lang(lang))


def ocr_image_to_html(source: Path, output: Path, lang: str = "ar+eng"):
    text=_ocr_text_file(source,lang); _write_text(output,"<!doctype html><html><body><pre>"+html.escape(text)+"</pre></body></html>\n")


def ocr_image_to_markdown(source: Path, output: Path, lang: str = "ar+eng"):
    text=_ocr_text_file(source,lang).strip(); _write_text(output,"# OCR Result\n\n"+text+"\n")


def ocr_pdf_to_markdown(source: Path, output: Path, lang: str = "ar+eng"):
    import pytesseract
    pdf=_pdf_doc(source); pages=[]
    try:
        for i,p in enumerate(pdf,1):
            pix=p.get_pixmap(dpi=160,alpha=False); img=Image.open(io.BytesIO(pix.tobytes("png"))); text=pytesseract.image_to_string(img,lang=_ocr_lang(lang))
            pages.append(f"## Page {i}\n\n{text.strip()}")
    finally: pdf.close()
    _write_text(output,"\n\n".join(pages)+"\n")


def ocr_pdf_to_csv(source: Path, output: Path, lang: str = "ar+eng"):
    import csv, pytesseract
    pdf=_pdf_doc(source); rows=["page,text"]
    try:
        for i,p in enumerate(pdf,1):
            pix=p.get_pixmap(dpi=150,alpha=False); img=Image.open(io.BytesIO(pix.tobytes("png"))); text=pytesseract.image_to_string(img,lang=_ocr_lang(lang)).replace('"','""').replace("\n"," ").strip(); rows.append(f'{i},"{text}"')
    finally: pdf.close()
    _write_text(output,"\n".join(rows)+"\n")


def ocr_image_to_csv(source: Path, output: Path, lang: str = "ar+eng"):
    import csv
    text=_ocr_text_file(source,lang)
    rows=[r for r in text.splitlines() if r.strip()]
    with output.open("w",encoding="utf-8",newline="") as f:
        w=csv.writer(f); w.writerow(["line","text"]); w.writerows((i+1,r) for i,r in enumerate(rows))


def ocr_receipt_fields(source: Path, output: Path, lang: str = "ar+eng"):
    text=_ocr_text_file(source,lang); patterns={"emails":r"[\w.+-]+@[\w-]+\.[\w.-]+","phones":r"(?:\+?\d[\d\s().-]{6,}\d)","dates":r"\b\d{1,4}[/-]\d{1,2}[/-]\d{1,4}\b","money":r"(?:[$€£﷼]|SAR|USD|EUR)\s?\d+(?:[.,]\d+)?"}
    out={k:re.findall(v,text,re.I) for k,v in patterns.items()}; _write_text(output,json.dumps(out,ensure_ascii=False,indent=2))


def ocr_invoice_fields(source: Path, output: Path, lang: str = "ar+eng"):
    text=_ocr_text_file(source,lang)
    total=None
    matches=re.findall(r"(?:total|الإجمالي|grand total)\D{0,20}([\d,.]+)",text,re.I)
    if matches: total=matches[-1]
    _write_text(output,json.dumps({"invoice_number":(re.findall(r"(?:invoice|فاتورة)\D{0,20}([A-Z0-9-]+)",text,re.I) or [None])[0],"total":total,"raw":text},ensure_ascii=False,indent=2))


def ocr_deduplicate(source: Path, output: Path, lang: str = "ar+eng"):
    text=_ocr_text_file(source,lang); seen=set(); lines=[]
    for line in text.splitlines():
        key=re.sub(r"\s+"," ",line).strip().casefold()
        if key and key not in seen: seen.add(key); lines.append(line.strip())
    _write_text(output,"\n".join(lines)+"\n")


def ocr_entities(source: Path, output: Path, lang: str = "ar+eng"):
    text=_ocr_text_file(source,lang); patterns={"emails":r"[\w.+-]+@[\w-]+\.[\w.-]+","urls":r"https?://\S+","phones":r"(?:\+?\d[\d\s().-]{6,}\d)","dates":r"\b\d{1,4}[/-]\d{1,2}[/-]\d{1,4}\b"}
    _write_text(output,json.dumps({k:sorted(set(re.findall(p,text,re.I))) for k,p in patterns.items()},ensure_ascii=False,indent=2))


def ocr_language_report(source: Path, output: Path, lang: str = "ar+eng"):
    text=_ocr_text_file(source,lang); ar=len(re.findall(r"[\u0600-\u06ff]",text)); en=len(re.findall(r"[A-Za-z]",text)); total=max(1,ar+en)
    _write_text(output,json.dumps({"arabic_ratio":round(ar/total,3),"english_ratio":round(en/total,3),"detected":"ar+en" if ar and en else "ar" if ar else "en" if en else "unknown"},ensure_ascii=False,indent=2))


def bzip2_compress(source: Path, output: Path):
    output.write_bytes(bz2.compress(source.read_bytes(),compresslevel=9))

def bzip2_decompress(source: Path, output: Path):
    try: output.write_bytes(bz2.decompress(source.read_bytes()))
    except OSError as exc: raise ValueError("ملف BZIP2 غير صالح.") from exc

def xz_compress(source: Path, output: Path): output.write_bytes(lzma.compress(source.read_bytes(),preset=9))
def xz_decompress(source: Path, output: Path):
    try: output.write_bytes(lzma.decompress(source.read_bytes()))
    except lzma.LZMAError as exc: raise ValueError("ملف XZ غير صالح.") from exc

def tar_gzip_create(files, output: Path):
    with tarfile.open(output,"w:gz") as tf:
        for path,name in files: tf.add(path,arcname=name,recursive=False)

def tar_gzip_extract(source: Path, output_dir: Path):
    output_dir.mkdir(parents=True,exist_ok=True); out=[]
    with tarfile.open(source,"r:gz") as tf:
        for m in tf.getmembers():
            if not (m.isfile() or m.isdir()) or Path(m.name).is_absolute() or ".." in Path(m.name).parts: raise ValueError("الأرشيف يحتوي على مسار غير آمن.")
        for m in tf.getmembers():
            target=output_dir/m.name
            target.parent.mkdir(parents=True,exist_ok=True); tf.extract(m,output_dir,set_attrs=False); out.append(target)
    return out

def zip_duplicate_report(source: Path, output: Path):
    with zipfile.ZipFile(source) as zf:
        names=[i.filename for i in zf.infolist()]; seen=set(); dup=[]
        for n in names:
            if n in seen: dup.append(n)
            seen.add(n)
        _write_text(output,json.dumps({"duplicates":sorted(set(dup)),"entries":len(names)},ensure_ascii=False,indent=2))

def tar_integrity(source: Path, output: Path):
    try:
        with tarfile.open(source,"r:*") as tf: members=tf.getmembers(); ok=all(not Path(m.name).is_absolute() and ".." not in Path(m.name).parts for m in members)
    except tarfile.TarError: ok=False; members=[]
    _write_text(output,json.dumps({"valid":ok,"entries":len(members)},ensure_ascii=False,indent=2))

def archive_size_report(source: Path, output: Path):
    ext=source.suffix.lower(); raw=source.stat().st_size; unpacked=None; entries=None
    try:
        if ext==".zip":
            with zipfile.ZipFile(source) as zf: infos=zf.infolist(); unpacked=sum(i.file_size for i in infos); entries=len(infos)
        elif ext in {".tar",".gz",".tgz",".bz2",".xz"}:
            with tarfile.open(source,"r:*") as tf: ms=tf.getmembers(); unpacked=sum(m.size for m in ms if m.isfile()); entries=len(ms)
    except Exception: pass
    _write_text(output,json.dumps({"compressed_bytes":raw,"uncompressed_bytes":unpacked,"entries":entries,"ratio":round(unpacked/raw,3) if unpacked is not None and raw else None},ensure_ascii=False,indent=2))

def tar_bzip2_create(files, output: Path):
    with tarfile.open(output,"w:bz2") as tf:
        for path,name in files: tf.add(path,arcname=name,recursive=False)

def tar_bzip2_extract(source: Path, output_dir: Path):
    output_dir.mkdir(parents=True,exist_ok=True); out=[]
    with tarfile.open(source,"r:bz2") as tf:
        for m in tf.getmembers():
            if Path(m.name).is_absolute() or ".." in Path(m.name).parts: raise ValueError("الأرشيف يحتوي على مسار غير آمن.")
        for m in tf.getmembers(): tf.extract(m,output_dir,set_attrs=False); out.append(output_dir/m.name)
    return out

def base64_decode(source: Path, output: Path):
    try: output.write_bytes(base64.b64decode(source.read_text(encoding="utf-8"),validate=True))
    except Exception as exc: raise ValueError("Base64 غير صالح.") from exc

def url_encode(source: Path, output: Path): _write_text(output,urllib.parse.quote(source.read_text(encoding="utf-8")))
def url_decode(source: Path, output: Path): _write_text(output,urllib.parse.unquote(source.read_text(encoding="utf-8")))
def hex_encode(source: Path, output: Path): _write_text(output, source.read_bytes().hex())
def json_minify(source: Path, output: Path): _write_text(output,json.dumps(json.loads(source.read_text(encoding="utf-8")),ensure_ascii=False,separators=(",",":")))
def text_diff(a: Path,b: Path,output: Path):
    diff=list(difflib.unified_diff(a.read_text(encoding="utf-8").splitlines(),b.read_text(encoding="utf-8").splitlines(),fromfile=a.name,tofile=b.name,lineterm=""))
    text="\n".join(diff).strip()
    _write_text(output,(text+"\n") if text else "No differences found. The files are identical.\n")
def checksum_compare(a: Path,b: Path,output: Path):
    def h(p): return hashlib.sha256(p.read_bytes()).hexdigest()
    ha,hb=h(a),h(b); _write_text(output,json.dumps({"same":ha==hb,"sha256_a":ha,"sha256_b":hb},indent=2))
def uuid_list(count: str, output: Path):
    n=max(1,min(500,int(count or 10))); _write_text(output,"\n".join(str(uuid.uuid4()) for _ in range(n))+"\n")
def regex_extract(source: Path, output: Path, pattern: str):
    try: rx=re.compile(pattern or r"\b\w+\b",re.I|re.UNICODE)
    except re.error as exc: raise ValueError("التعبير المنتظم غير صالح.") from exc
    _write_text(output,"\n".join(sorted(set(rx.findall(source.read_text(encoding="utf-8")))))+"\n")

def rename_extension_report(source: Path, output: Path):
    p=Path(source.name); _write_text(output,json.dumps({"filename":p.name,"extension":p.suffix.lower(),"stem":p.stem,"mime":mimetypes.guess_type(p.name)[0]},ensure_ascii=False,indent=2))

# IDs handled by the mega dispatcher.
PDF_IDS={"pdf-to-docx","pdf-to-markdown","pdf-compare","pdf-repair","pdf-image-extract","pdf-links-report","pdf-annotations-report","pdf-page-size-report","pdf-redact","pdf-unlock"}
IMAGE_IDS={"image-upscale","image-blur","image-pixelate","image-invert","image-posterize","image-color-palette","image-watermark","image-background-cleaner","image-auto-orient","image-round-corners"}
OFFICE_IDS={"docx-to-markdown","docx-table-to-csv","xlsx-to-html","xlsx-summary","csv-to-markdown","csv-statistics","json-to-html","html-to-text","markdown-to-text","pptx-to-markdown"}
OCR_IDS={"ocr-image-to-html","ocr-image-to-markdown","ocr-pdf-to-markdown","ocr-pdf-to-csv","ocr-image-to-csv","ocr-receipt-fields","ocr-invoice-fields","ocr-text-deduplicate","ocr-entities","ocr-language-report"}
ARCHIVE_IDS={"bzip2-compress","bzip2-decompress","xz-compress","xz-decompress","tar-gzip-create","tar-gzip-extract","zip-duplicate-report","tar-integrity","tar-bzip2-create","tar-bzip2-extract"}
UTILITY_IDS={"base64-decode","url-encode","url-decode","json-minify","text-diff","checksum-compare","uuid-list-generator","regex-extract","file-extension-report","hex-encode"}
COMBINE_IDS={"pdf-compare","text-diff","checksum-compare","tar-gzip-create","tar-bzip2-create"}
NO_INPUT_IDS={"uuid-list-generator"}
