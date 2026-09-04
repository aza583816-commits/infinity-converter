import csv
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from docx import Document
from openpyxl import load_workbook, Workbook
from pptx import Presentation


def docx_to_text(source: Path, output: Path):
    doc = Document(str(source))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append('\t'.join(cell.text.strip() for cell in row.cells))
    text = '\n'.join(parts).strip()
    if not text:
        raise ValueError('لم يتم العثور على نص في مستند Word.')
    output.write_text(text, encoding='utf-8')


def docx_to_html(source: Path, output: Path):
    doc = Document(str(source))
    chunks = ['<!doctype html><html><head><meta charset="utf-8"><title>Document</title></head><body>']
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            tag = 'h2' if p.style and p.style.name and 'Heading' in p.style.name else 'p'
            chunks.append(f'<{tag}>{_html_escape(text)}</{tag}>')
    for table in doc.tables:
        chunks.append('<table><tbody>')
        for row in table.rows:
            chunks.append('<tr>' + ''.join(f'<td>{_html_escape(c.text)}</td>' for c in row.cells) + '</tr>')
        chunks.append('</tbody></table>')
    chunks.append('</body></html>')
    output.write_text(''.join(chunks), encoding='utf-8')


def xlsx_to_csv(source: Path, output: Path):
    wb = load_workbook(source, read_only=True, data_only=True)
    try:
        if not wb.sheetnames:
            raise ValueError('ملف Excel لا يحتوي أوراقًا.')
        ws = wb[wb.sheetnames[0]]
        with output.open('w', newline='', encoding='utf-8') as fh:
            writer = csv.writer(fh)
            for row in ws.iter_rows(values_only=True):
                writer.writerow([v if v is not None else '' for v in row])
    finally:
        wb.close()


def xlsx_to_json(source: Path, output: Path):
    wb = load_workbook(source, read_only=True, data_only=True)
    try:
        result = {}
        for name in wb.sheetnames:
            ws = wb[name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                result[name] = []
                continue
            headers = [str(v).strip() if v is not None else f'column_{i+1}' for i, v in enumerate(rows[0])]
            data = []
            for row in rows[1:]:
                item = {headers[i]: (row[i] if i < len(row) else None) for i in range(len(headers))}
                data.append(item)
            result[name] = data
        output.write_text(json.dumps(result, ensure_ascii=False, indent=2, default=str), encoding='utf-8')
    finally:
        wb.close()


def pptx_to_text(source: Path, output: Path):
    prs = Presentation(str(source))
    parts = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_text.append(shape.text.strip())
        parts.append(f'--- Slide {index} ---\n' + '\n'.join(slide_text))
    output.write_text('\n\n'.join(parts), encoding='utf-8')


def csv_to_json(source: Path, output: Path):
    with source.open('r', encoding='utf-8-sig', newline='') as fh:
        reader = csv.DictReader(fh)
        if reader.fieldnames is None:
            raise ValueError('ملف CSV لا يحتوي رؤوس أعمدة.')
        rows = list(reader)
    output.write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding='utf-8')


def json_to_csv(source: Path, output: Path):
    data = json.loads(source.read_text(encoding='utf-8'))
    if not isinstance(data, list) or not all(isinstance(item, dict) for item in data):
        raise ValueError('JSON يجب أن يكون قائمة من كائنات متشابهة.')
    headers = []
    for item in data:
        for key in item:
            if key not in headers:
                headers.append(key)
    with output.open('w', encoding='utf-8', newline='') as fh:
        writer = csv.DictWriter(fh, fieldnames=headers)
        writer.writeheader()
        for item in data:
            writer.writerow({key: _scalar(value) for key, value in item.items()})


def json_to_xlsx(source: Path, output: Path):
    data = json.loads(source.read_text(encoding='utf-8'))
    if not isinstance(data, list) or not all(isinstance(item, dict) for item in data):
        raise ValueError('JSON يجب أن يكون قائمة من كائنات متشابهة.')
    headers = []
    for item in data:
        for key in item:
            if key not in headers:
                headers.append(key)
    wb = Workbook()
    ws = wb.active
    ws.append(headers)
    for item in data:
        ws.append([_scalar(item.get(key)) for key in headers])
    wb.save(output)


def xml_to_json(source: Path, output: Path):
    raw = source.read_bytes()
    upper = raw[:2 * 1024 * 1024].upper()
    if b'<!DOCTYPE' in upper or b'<!ENTITY' in upper:
        raise ValueError('محتوى XML الذي يحتوي على DOCTYPE أو ENTITY غير مدعوم لأسباب أمنية.')
    root = ET.fromstring(raw)
    def node_to_obj(node):
        children = list(node)
        payload = {}
        if node.attrib:
            payload['@attributes'] = dict(node.attrib)
        if children:
            for child in children:
                value = node_to_obj(child)
                key = child.tag
                if key in payload:
                    if not isinstance(payload[key], list):
                        payload[key] = [payload[key]]
                    payload[key].append(value)
                else:
                    payload[key] = value
            text = (node.text or '').strip()
            if text:
                payload['#text'] = text
            return payload
        text = (node.text or '').strip()
        if payload:
            if text:
                payload['#text'] = text
            return payload
        return text
    output.write_text(json.dumps({root.tag: node_to_obj(root)}, ensure_ascii=False, indent=2), encoding='utf-8')


def text_to_json(source: Path, output: Path):
    lines = [line.rstrip('\n') for line in source.read_text(encoding='utf-8').splitlines()]
    output.write_text(json.dumps({'lines': lines, 'line_count': len(lines)}, ensure_ascii=False, indent=2), encoding='utf-8')


def json_pretty(source: Path, output: Path):
    data = json.loads(source.read_text(encoding='utf-8'))
    output.write_text(json.dumps(data, ensure_ascii=False, indent=2, sort_keys=True), encoding='utf-8')


def _scalar(value):
    if isinstance(value, (dict, list)):
        return json.dumps(value, ensure_ascii=False)
    return '' if value is None else value


def _html_escape(value):
    import html
    return html.escape(value or '')
